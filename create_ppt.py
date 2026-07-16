from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    
    # --- Slide 1: Title ---
    slide1 = prs.slides.add_slide(title_slide_layout)
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    
    title1.text = "🍜 FoodieGo"
    subtitle1.text = "Hệ thống đặt món trực tuyến\nĐồ án môn SPQM | Level 2\nTrọng tâm: LEAN + Lead Time + Measurement"
    
    # --- Slide 2: Giới thiệu ---
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes2 = slide2.shapes
    title2 = shapes2.title
    body2 = shapes2.placeholders[1]
    
    title2.text = "📌 Giới thiệu dự án"
    tf2 = body2.text_frame
    tf2.text = "FoodieGo là nền tảng kết nối toàn diện với 3 vai trò chính:"
    
    p = tf2.add_paragraph()
    p.text = "👩‍🦰 Customer (Khách hàng): Xem nhà hàng, menu, đặt món và theo dõi đơn hàng."
    p.level = 1
    
    p = tf2.add_paragraph()
    p.text = "🏪 Restaurant Owner (Chủ nhà hàng): Quản lý thông tin, menu và xử lý đơn hàng."
    p.level = 1
    
    p = tf2.add_paragraph()
    p.text = "👨‍💻 Admin (Quản trị viên): Quản lý người dùng, nhà hàng và xem thống kê hệ thống."
    p.level = 1
    
    # --- Slide 3: Đội ngũ phát triển ---
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    shapes3 = slide3.shapes
    title3 = shapes3.title
    body3 = shapes3.placeholders[1]
    
    title3.text = "👥 Đội ngũ phát triển"
    tf3 = body3.text_frame
    tf3.text = "Thành viên nhóm và vai trò:"
    
    p = tf3.add_paragraph()
    p.text = "Cao Đình Bảo (Team Lead): Auth + Users + CI/CD"
    p.level = 1
    
    p = tf3.add_paragraph()
    p.text = "Võ Duy Hoàng (Backend): Restaurants + Menu Items"
    p.level = 1
    
    p = tf3.add_paragraph()
    p.text = "Phạm Hải Thiên (Backend): Orders + Delivery Service"
    p.level = 1
    
    # --- Slide 4: Tech Stack ---
    slide4 = prs.slides.add_slide(bullet_slide_layout)
    shapes4 = slide4.shapes
    title4 = shapes4.title
    body4 = shapes4.placeholders[1]
    
    title4.text = "🏗️ Công nghệ sử dụng"
    tf4 = body4.text_frame
    tf4.text = "Hệ thống được xây dựng trên các công nghệ hiện đại:"
    
    p = tf4.add_paragraph()
    p.text = "Core Backend: Node.js + Express"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "Database: MySQL"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "Microservice: Python FastAPI (Tính phí giao hàng)"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "Bảo mật: JWT + Role-based Authorization"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "Kiểm thử & Chất lượng: Jest, Supertest, ESLint, SonarQube"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "DevOps: Docker Compose, GitHub Actions (CI/CD)"
    p.level = 1
    
    # --- Slide 5: Quality Metrics ---
    slide5 = prs.slides.add_slide(bullet_slide_layout)
    shapes5 = slide5.shapes
    title5 = shapes5.title
    body5 = shapes5.placeholders[1]
    
    title5.text = "📊 Mục tiêu chất lượng"
    tf5 = body5.text_frame
    tf5.text = "Dự án áp dụng quy trình kiểm soát chất lượng (DoD) nghiêm ngặt:"
    
    p = tf5.add_paragraph()
    p.text = "Test Coverage: ≥ 80%"
    p.level = 1
    p = tf5.add_paragraph()
    p.text = "Tỉ lệ lỗi CI (CI Fail Rate): ≤ 10%"
    p.level = 1
    p = tf5.add_paragraph()
    p.text = "Thời gian review (Lead Time): ≤ 1 ngày"
    p.level = 1
    p = tf5.add_paragraph()
    p.text = "Lỗi SonarQube: ≤ 10 issues"
    p.level = 1
    
    # --- Slide 6: Kết luận ---
    slide6 = prs.slides.add_slide(title_slide_layout)
    title6 = slide6.shapes.title
    subtitle6 = slide6.placeholders[1]
    
    title6.text = "🎉 Cảm ơn các bạn đã lắng nghe!"
    subtitle6.text = "Q&A\n\nFoodieGo © 2024 — Đồ án SPQM"
    
    # Save presentation
    prs.save("FoodieGo_Presentation.pptx")
    print("Presentation saved as FoodieGo_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
